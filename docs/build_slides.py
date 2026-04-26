"""Build the Energy Hackathon presentation deck.

Generates ``docs/energy-hackathon-deck.pptx`` plus the matplotlib figures it
embeds (under ``docs/architecture/``). The architecture slide is rendered as
native PowerPoint shapes so it can be tweaked inside PowerPoint after the
fact; the frame-snap timeline is a PNG because it has too many tiny scored
markers to be worth doing by hand.

Run with::

    python3 docs/build_slides.py

Requires ``python-pptx`` and ``matplotlib`` on the local interpreter.
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paths and palette
# ---------------------------------------------------------------------------

DOCS = Path(__file__).resolve().parent
ASSETS = DOCS / "architecture"
ASSETS.mkdir(parents=True, exist_ok=True)
ICONS = DOCS / "icons" / "aws"
OUT = DOCS / "energy-hackathon-deck.pptx"


def _icon(name: str) -> Path:
    """Resolve an AWS icon by short name. Prefer the @5x PNGs for crisp rendering."""
    base = ICONS / f"Arch_{name}_64@5x.png"
    if base.exists():
        return base
    return ICONS / f"Arch_{name}_64.png"


# Aliases for the icons we use, so callers don't have to remember filenames.
AWS_ICONS = {
    "alb": _icon("Elastic-Load-Balancing"),
    "ecs": _icon("Amazon-Elastic-Container-Service"),
    "fargate": _icon("AWS-Fargate"),
    "ecr": _icon("Amazon-Elastic-Container-Registry"),
    "rds": _icon("Amazon-RDS"),
    "s3": _icon("Amazon-Simple-Storage-Service"),
    "bedrock": _icon("Amazon-Bedrock"),
    "secrets": _icon("AWS-Secrets-Manager"),
    "iam": _icon("AWS-Identity-and-Access-Management"),
    "logs": _icon("Amazon-CloudWatch"),
    "lambda": _icon("AWS-Lambda"),
    "events": _icon("Amazon-EventBridge"),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x4B, 0x55, 0x63)
SUBTLE = RGBColor(0x9C, 0xA3, 0xAF)
BG = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x0F, 0x76, 0x6E)
ACCENT_DARK = RGBColor(0x0B, 0x4D, 0x48)
HIGHLIGHT = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xDC, 0x26, 0x26)

STORAGE_FILL = RGBColor(0xDB, 0xEA, 0xFE)
STORAGE_LINE = RGBColor(0x1E, 0x3A, 0x8A)
COMPUTE_FILL = RGBColor(0xDC, 0xFC, 0xE7)
COMPUTE_LINE = RGBColor(0x14, 0x53, 0x2D)
AI_FILL = RGBColor(0xED, 0xE9, 0xFE)
AI_LINE = RGBColor(0x5B, 0x21, 0xB6)
SEC_FILL = RGBColor(0xFC, 0xE7, 0xF3)
SEC_LINE = RGBColor(0x9D, 0x17, 0x4B)
LOG_FILL = RGBColor(0xFE, 0xF3, 0xC7)
LOG_LINE = RGBColor(0x92, 0x40, 0x0E)
USER_FILL = RGBColor(0xFF, 0xED, 0xD5)
USER_LINE = RGBColor(0x9A, 0x34, 0x12)
VPC_FILL = RGBColor(0xF8, 0xFA, 0xFC)
VPC_LINE = RGBColor(0x94, 0xA3, 0xB8)

# AWS Solution Design palette — matches the sister project's
# `architecture_aws.svg` so this deck visually matches our other AWS docs.
AWS_CLOUD_LINE = RGBColor(0x23, 0x2F, 0x3E)   # navy
AWS_REGION_LINE = RGBColor(0x00, 0xA4, 0xA6)  # teal
AWS_VPC_LINE = RGBColor(0x8C, 0x4F, 0xFF)     # purple
AWS_SUBNET_LINE = RGBColor(0x7A, 0xA1, 0x16)  # green
AWS_PRIVATE_SUBNET_LINE = RGBColor(0x1F, 0x6F, 0xEB)  # blue (informational)
AWS_AVAIL_TAB_TEXT = BG  # tab labels are always white over a coloured fill


# ---------------------------------------------------------------------------
# Generic helpers
# ---------------------------------------------------------------------------


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Inter",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, *, fill, line, line_w=1.0, corner=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = corner
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_label(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    fill,
    line,
    title=None,
    size=12,
    bold=False,
    color=INK,
):
    box = add_box(slide, x, y, w, h, fill=fill, line=line, line_w=1.25)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(75000)
    tf.margin_right = Emu(75000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if title:
        r = p.add_run()
        r.text = title
        r.font.name = "Inter"
        r.font.size = Pt(size + 1)
        r.font.bold = True
        r.font.color.rgb = color
        if text:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = text
            r2.font.name = "Inter"
            r2.font.size = Pt(size - 1)
            r2.font.color.rgb = MUTED
    else:
        r = p.add_run()
        r.text = text
        r.font.name = "Inter"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def add_connector(
    slide,
    x1,
    y1,
    x2,
    y2,
    *,
    label=None,
    color=MUTED,
    weight=1.5,
    dashed=False,
    label_dy=Inches(-0.18),
    label_dx=Inches(0),
):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    ln = line.line._get_or_add_ln()
    if dashed:
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "dash")
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")

    if label:
        mx = int((x1 + x2) / 2)
        my = int((y1 + y2) / 2)
        lab_w = Inches(2.4)
        lab_h = Inches(0.3)
        tb = slide.shapes.add_textbox(mx - lab_w // 2 + label_dx, my + label_dy, lab_w, lab_h)
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Emu(20000)
        tf.margin_top = tf.margin_bottom = Emu(20000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Inter"
        r.font.size = Pt(10)
        r.font.color.rgb = MUTED
    return line


def add_aws_container(slide, x, y, w, h, *, color, label):
    """Draw the AWS-style dashed-border container with a coloured tab label.

    Mirrors the visual language of `mission-planning-demo/docs/architecture_aws.svg`
    so this deck reads like a sibling AWS Solution Design diagram.
    """
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.background()
    rect.line.color.rgb = color
    rect.line.width = Pt(1.25)
    ln = rect.line._get_or_add_ln()
    prst = etree.SubElement(ln, qn("a:prstDash"))
    prst.set("val", "dash")
    rect.shadow.inherit = False

    tab_w = Inches(max(0.95, 0.085 * len(label) + 0.25))
    tab_h = Inches(0.24)
    tab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, tab_w, tab_h)
    tab.fill.solid()
    tab.fill.fore_color.rgb = color
    tab.line.fill.background()
    tab.shadow.inherit = False

    tb = slide.shapes.add_textbox(x, y, tab_w, tab_h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(10000)
    tf.margin_bottom = Emu(10000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label
    r.font.name = "Inter"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = AWS_AVAIL_TAB_TEXT
    return rect


def add_aws_icon(
    slide,
    icon_key,
    cx,
    cy_top,
    *,
    size=Inches(0.62),
    title=None,
    subtitle=None,
    label_w=Inches(2.0),
):
    """Place an official AWS service icon centred on `cx` with optional labels below.

    `icon_key` is a key from AWS_ICONS. Returns (icon_left, icon_top,
    icon_bottom_y) so callers can wire connectors to the icon's edges.
    """
    icon_path = AWS_ICONS[icon_key]
    icon_left = cx - size // 2
    pic = slide.shapes.add_picture(str(icon_path), icon_left, cy_top, width=size, height=size)
    pic.shadow.inherit = False
    label_top = cy_top + size + Inches(0.04)
    if title:
        add_text(
            slide,
            cx - label_w // 2,
            label_top,
            label_w,
            Inches(0.22),
            title,
            size=10,
            bold=True,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        label_top = label_top + Inches(0.20)
    if subtitle:
        add_text(
            slide,
            cx - label_w // 2,
            label_top,
            label_w,
            Inches(0.22),
            subtitle,
            size=8,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
    return icon_left, cy_top, cy_top + size


def add_title_bar(slide, title, *, kicker=None):
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.55), Inches(0.08), Inches(0.6)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    if kicker:
        add_text(
            slide,
            Inches(0.6),
            Inches(0.5),
            Inches(12),
            Inches(0.3),
            kicker.upper(),
            size=11,
            bold=True,
            color=ACCENT,
        )
        add_text(
            slide,
            Inches(0.6),
            Inches(0.78),
            Inches(12.4),
            Inches(0.6),
            title,
            size=26,
            bold=True,
            color=INK,
        )
    else:
        add_text(
            slide,
            Inches(0.6),
            Inches(0.55),
            Inches(12.4),
            Inches(0.6),
            title,
            size=26,
            bold=True,
            color=INK,
        )


def add_footer(slide, idx, total):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.1),
        Inches(8),
        Inches(0.3),
        "Energy Infrastructure Health · Multimodal Video Search",
        size=10,
        color=SUBTLE,
    )
    add_text(
        slide,
        Inches(11.8),
        Inches(7.1),
        Inches(1.2),
        Inches(0.3),
        f"{idx} / {total}",
        size=10,
        color=SUBTLE,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------


def slide_title(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    add_text(
        s,
        Inches(0.8),
        Inches(1.6),
        Inches(11.7),
        Inches(0.5),
        "ENERGY HACKATHON · BEDROCK + TWELVELABS",
        size=14,
        bold=True,
        color=ACCENT,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(2.1),
        Inches(11.7),
        Inches(2.0),
        "Multimodal video search\nfor grid inspection footage",
        size=44,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(4.4),
        Inches(11.7),
        Inches(1.2),
        "Clip embeddings find the moment.\nFrame embeddings show you exactly where it is.",
        size=22,
        color=MUTED,
    )

    chip_rows = [
        [
            ("Bedrock · Marengo + Pegasus", AI_FILL, AI_LINE),
            ("YOLO-seg · ultralytics", AI_FILL, AI_LINE),
            ("RDS Postgres + pgvector (HNSW)", STORAGE_FILL, STORAGE_LINE),
        ],
        [
            ("ECS Fargate · FastAPI portal", COMPUTE_FILL, COMPUTE_LINE),
            ("EventBridge · 4 Lambdas", COMPUTE_FILL, COMPUTE_LINE),
            ("3 Fargate workers · frame · pegasus · yolo", LOG_FILL, LOG_LINE),
        ],
    ]
    cy = Inches(5.85)
    for row in chip_rows:
        cx = Inches(0.8)
        for label, fill, line in row:
            w = Inches(max(1.6, 0.13 * len(label) + 0.7))
            add_label(s, cx, cy, w, Inches(0.5), label, fill=fill, line=line, size=12, bold=True)
            cx = cx + w + Inches(0.18)
        cy = cy + Inches(0.65)

    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 2 — Problem
# ---------------------------------------------------------------------------


def slide_problem(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(
        s,
        "Hours of inspection footage. One question. Twelve seconds to find the answer.",
        kicker="The problem",
    )

    pains = [
        (
            "Filename search is useless",
            "Drone passes are named pipeline_001.mp4 — the interesting frame is buried 4 minutes in.",
        ),
        (
            "Manual scrubbing doesn't scale",
            "Inspectors review hours of identical-looking transmission line and pipeline footage every week.",
        ),
        (
            "Single-frame search misses motion",
            "A still snapshot loses arc flashes, smoke plumes, and the audio cues from a faulty transformer.",
        ),
    ]
    cx = Inches(0.6)
    cy = Inches(2.0)
    cw = Inches(4.0)
    ch = Inches(2.7)
    for title, body in pains:
        card = add_box(
            s,
            cx,
            cy,
            cw,
            ch,
            fill=RGBColor(0xF8, 0xFA, 0xFC),
            line=RGBColor(0xCB, 0xD5, 0xE1),
            line_w=1.0,
            corner=0.04,
        )
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(180000)
        tf.margin_right = Emu(180000)
        tf.margin_top = Emu(180000)
        tf.margin_bottom = Emu(180000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.name = "Inter"
        r.font.color.rgb = INK
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(14)
        r2.font.name = "Inter"
        r2.font.color.rgb = MUTED
        cx = cx + cw + Inches(0.15)

    add_text(
        s,
        Inches(0.6),
        Inches(5.0),
        Inches(12.1),
        Inches(1.6),
        "We want a search box that takes a sentence, an image, or both — "
        "and jumps straight to the right second of the right video.",
        size=20,
        bold=True,
        color=INK,
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 3 — Architecture (native shapes from the .tf files)
# ---------------------------------------------------------------------------


def slide_architecture(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(
        s,
        "Synchronous read path: search and upload portal",
        kicker="Architecture · 1 of 3",
    )

    # === Containers (drawn first so icons/labels render on top) ===
    AWS_X, AWS_Y = Inches(1.55), Inches(1.65)
    AWS_W, AWS_H = Inches(11.30), Inches(4.95)
    add_aws_container(s, AWS_X, AWS_Y, AWS_W, AWS_H, color=AWS_CLOUD_LINE, label="AWS Cloud")

    REG_X, REG_Y = Inches(1.75), Inches(1.95)
    REG_W, REG_H = Inches(11.00), Inches(4.55)
    add_aws_container(s, REG_X, REG_Y, REG_W, REG_H, color=AWS_REGION_LINE, label="us-east-1")

    VPC_X, VPC_Y = Inches(1.95), Inches(2.30)
    VPC_W, VPC_H = Inches(2.95), Inches(4.10)
    add_aws_container(s, VPC_X, VPC_Y, VPC_W, VPC_H, color=AWS_VPC_LINE, label="VPC 10.42.0.0/16")

    SUB_Y, SUB_H = Inches(2.65), Inches(3.65)
    SUB_A_X, SUB_A_W = Inches(2.05), Inches(2.75)
    add_aws_container(s, SUB_A_X, SUB_Y, SUB_A_W, SUB_H, color=AWS_SUBNET_LINE, label="Public Subnet AZ-a")

    # === Operator (outside AWS Cloud) ===
    USER_CX = Inches(0.85)
    USER_CY = Inches(3.55)
    add_text(
        s,
        Inches(0.4),
        USER_CY - Inches(0.45),
        Inches(1.20),
        Inches(0.32),
        "Operator",
        size=12,
        bold=True,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.4),
        USER_CY - Inches(0.18),
        Inches(1.20),
        Inches(0.3),
        "browser",
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.4),
        USER_CY + Inches(0.10),
        Inches(1.20),
        Inches(0.3),
        "shared-token cookie",
        size=8,
        color=SUBTLE,
        align=PP_ALIGN.CENTER,
    )

    # === Icon coordinates ===
    ICON_SIZE = Inches(0.62)
    SUB_A_CX = SUB_A_X + SUB_A_W // 2
    ALB_CX = SUB_A_CX
    ALB_CY = Inches(2.85)
    ECS_CX = SUB_A_CX
    ECS_CY = Inches(4.05)
    RDS_CX = SUB_A_CX
    RDS_CY = Inches(5.25)

    # Right side: regional services (S3, Bedrock, Secrets, ECR, CloudWatch, IAM)
    SVC_COL1 = Inches(6.10)
    SVC_COL2 = Inches(8.55)
    SVC_COL3 = Inches(11.10)
    SVC_ROW1 = Inches(2.85)
    SVC_ROW2 = Inches(5.05)

    # === Connectors (drawn before icons so endpoints sit under the artwork) ===
    # Operator → ALB (HTTPS)
    add_connector(
        s,
        Inches(1.45),
        USER_CY,
        ALB_CX - ICON_SIZE // 2,
        ALB_CY + ICON_SIZE // 2,
        color=ACCENT,
        weight=2.0,
        label="HTTPS",
        label_dy=Inches(-0.24),
    )
    # ALB → ECS (vertical, in-subnet)
    add_connector(
        s,
        ALB_CX,
        ALB_CY + ICON_SIZE,
        ECS_CX,
        ECS_CY,
        color=ACCENT,
        weight=2.0,
    )
    # ECS → RDS (cosine ANN, vertical in-subnet) — start arrow well below the
    # ECS labels and put "cosine ANN" to the right of the arrow midpoint so it
    # doesn't overlap the ECS subtitle.
    add_connector(
        s,
        ECS_CX,
        ECS_CY + ICON_SIZE + Inches(0.46),
        RDS_CX,
        RDS_CY,
        color=ACCENT,
        weight=2.0,
        label="cosine ANN",
        label_dx=Inches(0.50),
        label_dy=Inches(-0.05),
    )
    # ECS → S3 (presigned URLs) — horizontal right
    add_connector(
        s,
        ECS_CX + ICON_SIZE // 2,
        ECS_CY + ICON_SIZE // 2,
        SVC_COL1 - ICON_SIZE // 2,
        SVC_ROW1 + ICON_SIZE // 2,
        color=ACCENT,
        weight=1.4,
        label="presigned URLs",
        label_dy=Inches(-0.32),
    )
    # ECS → Bedrock (sync InvokeModel) — horizontal right, slight downward
    add_connector(
        s,
        ECS_CX + ICON_SIZE // 2,
        ECS_CY + ICON_SIZE // 2,
        SVC_COL2 - ICON_SIZE // 2,
        SVC_ROW1 + ICON_SIZE // 2,
        color=ACCENT,
        weight=1.4,
        label="sync InvokeModel",
        label_dy=Inches(0.18),
    )
    # ECS → Secrets Manager (dashed) — horizontal right, downward
    add_connector(
        s,
        ECS_CX + ICON_SIZE // 2,
        ECS_CY + ICON_SIZE - Inches(0.05),
        SVC_COL1 - ICON_SIZE // 2,
        SVC_ROW2 + ICON_SIZE // 2,
        color=MUTED,
        weight=1.0,
        dashed=True,
        label="GetSecretValue",
        label_dy=Inches(0.18),
    )

    # === Icons (drawn last so they cover the connector tails cleanly) ===
    add_aws_icon(
        s,
        "alb",
        ALB_CX,
        ALB_CY,
        size=ICON_SIZE,
        title="Application LB",
        subtitle="HTTP :80 · /health · /api/*",
    )
    add_aws_icon(
        s,
        "ecs",
        ECS_CX,
        ECS_CY,
        size=ICON_SIZE,
        title="ECS · Fargate",
        subtitle="FastAPI portal · 512 / 1024",
    )
    add_aws_icon(
        s,
        "rds",
        RDS_CX,
        RDS_CY,
        size=ICON_SIZE,
        title="RDS · pgvector",
        subtitle="db.t4g.large · HNSW-tuned",
    )
    add_aws_icon(
        s,
        "s3",
        SVC_COL1,
        SVC_ROW1,
        size=ICON_SIZE,
        title="S3 bucket",
        subtitle="raw-videos/ · embeddings/",
    )
    add_aws_icon(
        s,
        "bedrock",
        SVC_COL2,
        SVC_ROW1,
        size=ICON_SIZE,
        title="Bedrock",
        subtitle="Marengo Embed 3.0",
    )
    add_aws_icon(
        s,
        "ecr",
        SVC_COL3,
        SVC_ROW1,
        size=ICON_SIZE,
        title="ECR",
        subtitle="portal image",
    )
    add_aws_icon(
        s,
        "secrets",
        SVC_COL1,
        SVC_ROW2,
        size=ICON_SIZE,
        title="Secrets Mgr",
        subtitle="portal token · DB URL",
    )
    add_aws_icon(
        s,
        "logs",
        SVC_COL2,
        SVC_ROW2,
        size=ICON_SIZE,
        title="CloudWatch",
        subtitle="/ecs/<project> · 7d",
    )
    add_aws_icon(
        s,
        "iam",
        SVC_COL3,
        SVC_ROW2,
        size=ICON_SIZE,
        title="IAM",
        subtitle="task + execution roles",
    )

    add_text(
        s,
        Inches(0.5),
        Inches(6.65),
        Inches(12.4),
        Inches(0.4),
        "Operator browser also hits S3 directly with presigned PUT/GET URLs minted by ECS. "
        "The async ingest pipeline (S3 EventBridge → 3 Lambdas + Fargate frame worker) populates the "
        "embeddings table this slide queries — see next slide.",
        size=11,
        color=MUTED,
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 4 — Async write path (the embedding pipeline that lands rows in RDS)
# ---------------------------------------------------------------------------


def slide_async_pipeline(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(
        s,
        "Async embed pipeline: drop a video, get clip + frame embeddings",
        kicker="Architecture · 2 of 3 · Marengo",
    )

    ICON_SIZE = Inches(0.62)

    # === Containers ===
    AWS_X, AWS_Y = Inches(0.5), Inches(1.65)
    AWS_W, AWS_H = Inches(12.35), Inches(4.95)
    add_aws_container(s, AWS_X, AWS_Y, AWS_W, AWS_H, color=AWS_CLOUD_LINE, label="AWS Cloud")

    REG_X, REG_Y = Inches(0.70), Inches(1.95)
    REG_W, REG_H = Inches(11.95), Inches(4.55)
    add_aws_container(s, REG_X, REG_Y, REG_W, REG_H, color=AWS_REGION_LINE, label="us-east-1")

    # VPC wraps the in-VPC components only (start_clip_embed, finalize_clip_embed,
    # frame-embed-worker, RDS). It sits in the lower portion of the canvas.
    VPC_X, VPC_Y = Inches(2.65), Inches(4.05)
    VPC_W, VPC_H = Inches(9.85), Inches(2.40)
    add_aws_container(s, VPC_X, VPC_Y, VPC_W, VPC_H, color=AWS_VPC_LINE, label="VPC")

    # === Icon coordinates ===
    # Top row (regional services, outside VPC)
    TOP_Y = Inches(2.55)
    S3_CX = Inches(1.60)
    EB_CX = Inches(3.30)
    SFT_CX = Inches(5.00)        # start_frame_task (no VPC)
    BEDROCK_CX = Inches(11.65)   # regional service

    # Bottom row (in VPC)
    BOT_Y = Inches(4.55)
    SCE_CX = Inches(3.45)        # start_clip_embed
    FCE_CX = Inches(5.45)        # finalize_clip_embed
    FW_CX = Inches(8.20)         # frame-embed-worker (Fargate)
    RDS_CX = Inches(11.65)       # RDS

    icon_top = TOP_Y
    icon_bot = BOT_Y
    icon_top_mid = TOP_Y + ICON_SIZE // 2
    icon_bot_mid = BOT_Y + ICON_SIZE // 2

    # === Connectors (drawn before icons) ===
    # S3 → EventBridge: ObjectCreated bucket-level notification
    add_connector(
        s,
        S3_CX + ICON_SIZE // 2,
        icon_top_mid,
        EB_CX - ICON_SIZE // 2,
        icon_top_mid,
        color=ACCENT,
        weight=1.7,
        label="ObjectCreated",
        label_dy=Inches(-0.24),
    )

    # EB → start_frame_task (top row, same row)
    add_connector(
        s,
        EB_CX + ICON_SIZE // 2,
        icon_top_mid,
        SFT_CX - ICON_SIZE // 2,
        icon_top_mid,
        color=ACCENT,
        weight=1.5,
        label="video_uploaded",
        label_dy=Inches(-0.24),
    )

    # EB → start_clip_embed (top → bottom, fan-out target #2)
    # Drop label here; the EB → start_frame_task label already explains the rule.
    add_connector(
        s,
        EB_CX,
        icon_top + ICON_SIZE,
        SCE_CX,
        icon_bot,
        color=ACCENT,
        weight=1.5,
    )

    # EB → finalize_clip_embed (clip_output_ready rule fires on output.json) —
    # the EB tile sits up-and-left of finalize_clip_embed, so this arrow has
    # the most vertical room for a label and gets the rule callout. Pull the
    # label well below the arrow midpoint so it sits in the gap above the VPC
    # tab without colliding with the EB subtitle.
    add_connector(
        s,
        EB_CX + ICON_SIZE // 2,
        icon_top + ICON_SIZE,
        FCE_CX,
        icon_bot,
        color=HIGHLIGHT,
        weight=1.5,
        label="clip_output_ready",
        label_dy=Inches(0.20),
        label_dx=Inches(0.30),
    )

    # start_frame_task → frame-embed-worker
    add_connector(
        s,
        SFT_CX,
        icon_top + ICON_SIZE,
        FW_CX,
        icon_bot,
        color=ACCENT,
        weight=1.5,
        label="ecs.run_task",
        label_dy=Inches(-0.10),
    )

    # start_clip_embed → Bedrock async (bottom → top)
    add_connector(
        s,
        SCE_CX + ICON_SIZE // 2,
        icon_bot_mid,
        BEDROCK_CX - ICON_SIZE // 2,
        icon_top + Inches(0.20),
        color=ACCENT,
        weight=1.4,
        label="StartAsyncInvoke",
        label_dy=Inches(-0.42),
    )

    # frame-embed-worker → Bedrock (sync InvokeModel ×8)
    add_connector(
        s,
        FW_CX + ICON_SIZE // 2,
        icon_bot,
        BEDROCK_CX - ICON_SIZE // 2,
        icon_top + ICON_SIZE - Inches(0.10),
        color=ACCENT,
        weight=1.2,
        label="8× InvokeModel",
        label_dy=Inches(-0.32),
    )

    # frame-embed-worker → S3 (PUT thumbnails) — drop the label, just keep
    # a thin dashed line so it's clear where thumbnails come from.
    add_connector(
        s,
        FW_CX - Inches(0.20),
        icon_bot,
        S3_CX + Inches(0.10),
        icon_top + ICON_SIZE,
        color=MUTED,
        weight=1.0,
        dashed=True,
    )

    # finalize_clip_embed → RDS — long horizontal hop. Keep the label
    # squarely between frame-embed-worker and RDS so it doesn't land on the
    # frame-embed-worker icon.
    add_connector(
        s,
        FCE_CX + ICON_SIZE // 2,
        icon_bot_mid - Inches(0.05),
        RDS_CX - ICON_SIZE // 2,
        icon_bot_mid - Inches(0.05),
        color=ACCENT,
        weight=1.5,
        label="INSERT clip",
        label_dy=Inches(-0.20),
        label_dx=Inches(1.55),
    )

    # frame-embed-worker → RDS
    add_connector(
        s,
        FW_CX + ICON_SIZE // 2,
        icon_bot_mid + Inches(0.05),
        RDS_CX - ICON_SIZE // 2,
        icon_bot_mid + Inches(0.05),
        color=ACCENT,
        weight=1.5,
        label="INSERT frame",
        label_dy=Inches(0.18),
    )

    # === Icons ===
    add_aws_icon(s, "s3", S3_CX, TOP_Y, size=ICON_SIZE, title="S3", subtitle="raw-videos/ · output.json")
    add_aws_icon(s, "events", EB_CX, TOP_Y, size=ICON_SIZE, title="EventBridge", subtitle="2 rules · 4 fan-outs total")
    add_aws_icon(
        s,
        "lambda",
        SFT_CX,
        TOP_Y,
        size=ICON_SIZE,
        title="start_frame_task",
        subtitle="λ · no VPC · dispatcher",
    )
    add_aws_icon(
        s,
        "bedrock",
        BEDROCK_CX,
        TOP_Y,
        size=ICON_SIZE,
        title="Bedrock",
        subtitle="Marengo · async + sync",
    )

    add_aws_icon(
        s,
        "lambda",
        SCE_CX,
        BOT_Y,
        size=ICON_SIZE,
        title="start_clip_embed",
        subtitle="λ in VPC · pg8000",
    )
    add_aws_icon(
        s,
        "lambda",
        FCE_CX,
        BOT_Y,
        size=ICON_SIZE,
        title="finalize_clip_embed",
        subtitle="λ in VPC · L2-norm + INSERT",
    )
    add_aws_icon(
        s,
        "fargate",
        FW_CX,
        BOT_Y,
        size=ICON_SIZE,
        title="frame-embed-worker",
        subtitle="Fargate · ffmpeg + 8× Bedrock",
    )
    add_aws_icon(
        s,
        "rds",
        RDS_CX,
        BOT_Y,
        size=ICON_SIZE,
        title="RDS · pgvector",
        subtitle="kind='clip' + kind='frame'",
    )

    add_text(
        s,
        Inches(0.5),
        Inches(6.65),
        Inches(12.4),
        Inches(0.4),
        "Both Marengo branches fan out from one bucket-level S3 → EventBridge notification. "
        "A third fan-out — start_yolo_task → yolo-detect-worker — runs in parallel; see the next slide.",
        size=11,
        color=MUTED,
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 5 — Async enrichments (Pegasus + YOLO)
# ---------------------------------------------------------------------------


def slide_enrichments(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(
        s,
        "Async enrichments: per-clip text + per-frame instance segmentation",
        kicker="Architecture · 3 of 3 · Pegasus + YOLO",
    )

    ICON_SIZE = Inches(0.55)

    # === Containers ===
    AWS_X, AWS_Y = Inches(0.5), Inches(1.65)
    AWS_W, AWS_H = Inches(12.35), Inches(4.95)
    add_aws_container(s, AWS_X, AWS_Y, AWS_W, AWS_H, color=AWS_CLOUD_LINE, label="AWS Cloud")

    REG_X, REG_Y = Inches(0.70), Inches(1.95)
    REG_W, REG_H = Inches(11.95), Inches(4.55)
    add_aws_container(s, REG_X, REG_Y, REG_W, REG_H, color=AWS_REGION_LINE, label="us-east-1")

    # === Coordinates ===
    LEFT_CY = Inches(4.05)         # S3 + EB vertical center
    R1_CY = Inches(2.55)           # Pegasus pipeline row
    R2_CY = Inches(5.45)           # YOLO pipeline row

    S3_CX = Inches(1.05)
    EB_CX = Inches(2.30)
    L_CX = Inches(3.85)            # finalize_clip_embed / start_yolo_task
    W_CX = Inches(5.65)            # workers
    BEDROCK_CX = Inches(7.55)      # Bedrock (Pegasus stream, only row 1)
    RDS_CX = Inches(11.40)
    RDS_CY = Inches(4.05)

    half = ICON_SIZE // 2

    # === Connectors (under icons) ===

    # 1. S3 → EB (output.json arrival fires clip_output_ready)
    add_connector(
        s,
        S3_CX + half,
        LEFT_CY + half,
        EB_CX - half,
        LEFT_CY + half,
        color=ACCENT,
        weight=1.6,
        label="ObjectCreated",
        label_dy=Inches(-0.24),
    )

    # 2. EB → finalize_clip_embed (clip_output_ready) — diagonal up-right.
    # Label sits below the line, in the gap between row 1 labels and EB icon.
    add_connector(
        s,
        EB_CX + half,
        LEFT_CY + half - Inches(0.05),
        L_CX - half,
        R1_CY + half,
        color=HIGHLIGHT,
        weight=1.5,
        label="clip_output_ready",
        label_dy=Inches(0.18),
        label_dx=Inches(-0.10),
    )

    # 3. EB → start_yolo_task (video_uploaded · 3rd fan-out) — diagonal down-right.
    add_connector(
        s,
        EB_CX + half,
        LEFT_CY + half + Inches(0.05),
        L_CX - half,
        R2_CY + half,
        color=ACCENT,
        weight=1.5,
        label="video_uploaded · 3rd fan-out",
        label_dy=Inches(0.20),
        label_dx=Inches(0.20),
    )

    # 4. finalize_clip_embed → clip-pegasus-worker
    add_connector(
        s,
        L_CX + half,
        R1_CY + half,
        W_CX - half,
        R1_CY + half,
        color=ACCENT,
        weight=1.4,
        label="ecs.run_task",
        label_dy=Inches(-0.22),
    )

    # 5. clip-pegasus-worker → Bedrock (Pegasus stream)
    add_connector(
        s,
        W_CX + half,
        R1_CY + half,
        BEDROCK_CX - half,
        R1_CY + half,
        color=ACCENT,
        weight=1.4,
        label="Pegasus stream",
        label_dy=Inches(-0.22),
    )

    # 6. clip-pegasus-worker → RDS (long diagonal down-right)
    # Routed below the Bedrock row so the "Pegasus stream" label stays clean.
    add_connector(
        s,
        W_CX + half,
        R1_CY + ICON_SIZE,
        RDS_CX - half,
        RDS_CY + half - Inches(0.10),
        color=ACCENT,
        weight=1.5,
        label="INSERT clip_descriptions",
        label_dy=Inches(0.10),
        label_dx=Inches(0.20),
    )

    # 7. start_yolo_task → yolo-detect-worker
    add_connector(
        s,
        L_CX + half,
        R2_CY + half,
        W_CX - half,
        R2_CY + half,
        color=ACCENT,
        weight=1.4,
        label="ecs.run_task",
        label_dy=Inches(-0.22),
    )

    # 8. yolo-detect-worker → S3 (read weights, dashed back-link)
    add_connector(
        s,
        W_CX - Inches(0.20),
        R2_CY,
        S3_CX + Inches(0.10),
        LEFT_CY + ICON_SIZE,
        color=MUTED,
        weight=1.0,
        dashed=True,
    )

    # 9. yolo-detect-worker → RDS (long diagonal up-right)
    add_connector(
        s,
        W_CX + half,
        R2_CY + half,
        RDS_CX - half,
        RDS_CY + half + Inches(0.10),
        color=ACCENT,
        weight=1.5,
        label="INSERT frame_detections",
        label_dy=Inches(-0.22),
        label_dx=Inches(0.40),
    )

    # === Icons ===
    add_aws_icon(
        s, "s3", S3_CX, LEFT_CY, size=ICON_SIZE,
        title="S3",
        subtitle="output.json · derived clips · YOLO weights",
    )
    add_aws_icon(
        s, "events", EB_CX, LEFT_CY, size=ICON_SIZE,
        title="EventBridge",
        subtitle="clip_output_ready · video_uploaded",
    )

    # Row 1 — Pegasus
    add_aws_icon(
        s, "lambda", L_CX, R1_CY, size=ICON_SIZE,
        title="finalize_clip_embed",
        subtitle="λ in VPC · INSERT clip + spawn",
    )
    add_aws_icon(
        s, "fargate", W_CX, R1_CY, size=ICON_SIZE,
        title="clip-pegasus-worker",
        subtitle="Fargate · ffmpeg cuts → derived/clips/",
    )
    add_aws_icon(
        s, "bedrock", BEDROCK_CX, R1_CY, size=ICON_SIZE,
        title="Bedrock · Pegasus 1.2",
        subtitle="invoke_model_with_response_stream",
    )

    # Row 2 — YOLO
    add_aws_icon(
        s, "lambda", L_CX, R2_CY, size=ICON_SIZE,
        title="start_yolo_task",
        subtitle="λ no VPC · ecs.run_task dispatcher",
    )
    add_aws_icon(
        s, "fargate", W_CX, R2_CY, size=ICON_SIZE,
        title="yolo-detect-worker",
        subtitle="Fargate · ultralytics CPU torch · polls frames",
    )

    # RDS shared
    add_aws_icon(
        s, "rds", RDS_CX, RDS_CY, size=ICON_SIZE,
        title="RDS · pgvector",
        subtitle="clip_descriptions · frame_detections",
    )

    add_text(
        s,
        Inches(0.5),
        Inches(6.65),
        Inches(12.4),
        Inches(0.4),
        "Pegasus runs once Marengo's clip rows are written; the worker cuts each clip, streams a "
        "natural-language description, and writes clip_descriptions. YOLO runs in parallel with the "
        "frame-embed worker and writes frame_detections that the UI paints as SVG polygons over thumbnails.",
        size=11,
        color=MUTED,
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 6 — Two granularities, one vector space
# ---------------------------------------------------------------------------


def slide_two_granularities(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(s, "Two granularities, one vector space", kicker="The model trick")

    LX, LY = Inches(0.6), Inches(1.85)
    LW, LH = Inches(6.0), Inches(3.85)
    add_box(
        s,
        LX,
        LY,
        LW,
        LH,
        fill=RGBColor(0xEF, 0xF6, 0xFF),
        line=STORAGE_LINE,
        line_w=1.5,
        corner=0.04,
    )
    add_text(
        s,
        LX + Inches(0.3),
        LY + Inches(0.2),
        LW - Inches(0.6),
        Inches(0.4),
        "CLIP",
        size=12,
        bold=True,
        color=STORAGE_LINE,
    )
    add_text(
        s,
        LX + Inches(0.3),
        LY + Inches(0.55),
        LW - Inches(0.6),
        Inches(0.7),
        "~6-second segments of the source video",
        size=20,
        bold=True,
    )
    bullets = [
        "bedrock:StartAsyncInvoke → output.json on S3",
        "1 row per (s3_key, segment_index, embeddingOption)",
        "embeddingOption ∈ { visual, audio, transcription }",
        "captures motion, audio, and speech in one shot",
        "cheap: 1 async job per video, runs out of band",
    ]
    by = LY + Inches(1.25)
    for b in bullets:
        add_text(
            s,
            LX + Inches(0.3),
            by,
            LW - Inches(0.6),
            Inches(0.4),
            "•  " + b,
            size=14,
            color=INK,
        )
        by = by + Inches(0.4)
    add_text(
        s,
        LX + Inches(0.3),
        LY + LH - Inches(0.5),
        LW - Inches(0.6),
        Inches(0.4),
        "kind = 'clip'   ·   timestamp_sec = (start_sec + end_sec) / 2",
        size=11,
        color=MUTED,
        font="JetBrains Mono",
    )

    RX = LX + LW + Inches(0.3)
    RY, RW, RH = LY, LW, LH
    add_box(
        s,
        RX,
        RY,
        RW,
        RH,
        fill=RGBColor(0xFE, 0xF7, 0xED),
        line=LOG_LINE,
        line_w=1.5,
        corner=0.04,
    )
    add_text(
        s,
        RX + Inches(0.3),
        RY + Inches(0.2),
        RW - Inches(0.6),
        Inches(0.4),
        "FRAME",
        size=12,
        bold=True,
        color=LOG_LINE,
    )
    add_text(
        s,
        RX + Inches(0.3),
        RY + Inches(0.55),
        RW - Inches(0.6),
        Inches(0.7),
        "1-fps still frames extracted with ffmpeg",
        size=20,
        bold=True,
    )
    bullets = [
        "bedrock:InvokeModel sync image API per frame",
        "1 row per (s3_key, frame_index)",
        "thumb_s3_key points at a 720px JPEG on S3",
        "frame-precise — one row per second of footage",
        "expensive but pinpoint: the thumbnail is exact",
    ]
    by = RY + Inches(1.25)
    for b in bullets:
        add_text(
            s,
            RX + Inches(0.3),
            by,
            RW - Inches(0.6),
            Inches(0.4),
            "•  " + b,
            size=14,
            color=INK,
        )
        by = by + Inches(0.4)
    add_text(
        s,
        RX + Inches(0.3),
        RY + RH - Inches(0.5),
        RW - Inches(0.6),
        Inches(0.4),
        "kind = 'frame'   ·   thumb_s3_key NOT NULL",
        size=11,
        color=MUTED,
        font="JetBrains Mono",
    )

    add_text(
        s,
        Inches(0.6),
        Inches(6.0),
        Inches(12.1),
        Inches(0.5),
        "Both kinds are 512-d L2-normalized Marengo vectors and share the same HNSW cosine index.",
        size=16,
        bold=True,
        color=ACCENT_DARK,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.6),
        Inches(6.45),
        Inches(12.1),
        Inches(0.4),
        "ORDER BY embedding <=> :query — clips and frames compete on equal footing. The result set is a mix of the two.",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        font="JetBrains Mono",
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 5 — Frame-snap timeline (matplotlib figure)
# ---------------------------------------------------------------------------


def render_frame_snap_figure(out_path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12.0, 4.0), dpi=200)
    ax.set_xlim(11.3, 18.7)
    ax.set_ylim(-0.4, 3.4)
    ax.axis("off")

    ax.annotate(
        "",
        xy=(18.7, 0),
        xytext=(11.3, 0),
        arrowprops=dict(arrowstyle="->", lw=1.4, color="#475569"),
    )
    for tick in range(12, 19):
        ax.plot([tick, tick], [-0.05, 0.05], color="#475569", lw=1)
        ax.text(tick, -0.18, f"{tick}s", color="#475569", ha="center", va="top", fontsize=9)

    clip_box = FancyBboxPatch(
        (12, 1.35),
        6.0,
        0.55,
        boxstyle="round,pad=0,rounding_size=0.07",
        fc="#DBEAFE",
        ec="#1E3A8A",
        lw=1.6,
    )
    ax.add_patch(clip_box)
    ax.text(
        15,
        1.625,
        "clip · segment_index = 2 · cosine score 0.74",
        ha="center",
        va="center",
        fontsize=11,
        color="#1E3A8A",
        fontweight="bold",
    )
    ax.text(12.0, 2.05, "start_sec = 12.0", ha="left", va="bottom", fontsize=9, color="#1E3A8A")
    ax.text(18.0, 2.05, "end_sec = 18.0", ha="right", va="bottom", fontsize=9, color="#1E3A8A")

    frame_data = [
        (12.0, 0.41, False),
        (13.0, 0.55, False),
        (14.0, 0.68, False),
        (15.0, 0.79, False),
        (15.3, 0.81, True),
        (16.0, 0.72, False),
        (17.0, 0.49, False),
        (18.0, 0.33, False),
    ]
    for ts, score, best in frame_data:
        if best:
            fc, ec, txt_color, marker_size = "#FDE68A", "#92400E", "#92400E", 280
        else:
            fc, ec, txt_color, marker_size = "#F1F5F9", "#94A3B8", "#475569", 150
        ax.scatter(
            [ts],
            [0.55],
            s=marker_size,
            marker="s",
            facecolor=fc,
            edgecolor=ec,
            lw=1.5,
            zorder=4,
        )
        ax.text(
            ts,
            0.92,
            f"{score:.2f}",
            ha="center",
            va="bottom",
            fontsize=9,
            color=txt_color,
            fontweight="bold" if best else "normal",
        )
    ax.text(11.7, 0.55, "frames\n(1 fps)", ha="right", va="center", fontsize=10, color="#475569")
    ax.text(11.7, 1.625, "clip\n(~6s)", ha="right", va="center", fontsize=10, color="#1E3A8A")

    ax.annotate(
        "",
        xy=(15.3, 0.78),
        xytext=(12.0, 1.35),
        arrowprops=dict(
            arrowstyle="->",
            lw=2.4,
            color="#F59E0B",
            connectionstyle="arc3,rad=-0.28",
        ),
    )
    ax.text(
        13.4,
        1.18,
        "snap thumbnail + #t = 15.30",
        ha="center",
        va="center",
        fontsize=11,
        color="#92400E",
        fontweight="bold",
    )

    ax.text(
        15.0,
        2.85,
        "result for this hit:\n"
        "score 0.74 (clip)  ·  timestamp_sec = 15.30  ·  thumb_s3_key = frames/.../frame_15.jpg",
        ha="center",
        va="center",
        fontsize=10,
        color="#0F172A",
        bbox=dict(boxstyle="round,pad=0.4", fc="#FEF3C7", ec="#92400E", lw=1.0),
    )

    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def slide_frame_snap(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(
        s,
        "The frame-snap trick: clips for recall, frames for the seek bar",
        kicker="Search refinement",
    )

    fig_path = ASSETS / "frame_snap.png"
    render_frame_snap_figure(fig_path)
    s.shapes.add_picture(str(fig_path), Inches(0.6), Inches(1.75), width=Inches(12.1))

    add_text(
        s,
        Inches(0.6),
        Inches(5.85),
        Inches(12.1),
        Inches(0.5),
        "Per clip in the candidate pool, snap timestamp_sec to the highest-scoring frame whose timestamp falls in [start_sec, end_sec].",
        size=14,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        Inches(0.6),
        Inches(6.3),
        Inches(12.1),
        Inches(0.5),
        "Use that frame's S3 thumbnail and its timestamp in the presigned URL fragment (#t=<sec>). "
        "Implementation: app/search.py::_refine_and_dedupe.",
        size=12,
        color=MUTED,
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 6 — Search algorithm (SQL + Python)
# ---------------------------------------------------------------------------


def slide_search_algo(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(s, "What runs on every /api/search/* call", kicker="The algorithm")

    LX, LY = Inches(0.6), Inches(1.85)
    LW, LH = Inches(6.2), Inches(4.7)
    sql_box = add_box(
        s,
        LX,
        LY,
        LW,
        LH,
        fill=RGBColor(0x0F, 0x17, 0x2A),
        line=RGBColor(0x0F, 0x17, 0x2A),
        corner=0.02,
    )
    tf = sql_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(220000)
    tf.margin_right = Emu(220000)
    tf.margin_top = Emu(180000)
    tf.margin_bottom = Emu(180000)
    code_lines = [
        ("-- 1. Pull a candidate pool, ordered by cosine ANN", "94A3B8"),
        ("SELECT", "7DD3FC"),
        ("    s3_key, kind, segment_index, frame_index,", "E2E8F0"),
        ("    start_sec, end_sec, timestamp_sec, thumb_s3_key,", "E2E8F0"),
        ("    1 - (embedding <=> :q::vector) AS score", "E2E8F0"),
        ("FROM embeddings", "7DD3FC"),
        ("ORDER BY embedding <=> :q::vector", "7DD3FC"),
        ("LIMIT 80;   -- = max(80, top_k * 4)", "7DD3FC"),
        ("", "E2E8F0"),
        ("-- HNSW index on embeddings(embedding vector_cosine_ops)", "94A3B8"),
        ("-- WITH (m = 16, ef_construction = 64)", "94A3B8"),
    ]
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for i, (text, color_hex) in enumerate(code_lines):
        para = p if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        r = para.add_run()
        r.text = text if text else " "
        r.font.name = "JetBrains Mono"
        r.font.size = Pt(13)
        r.font.color.rgb = RGBColor.from_string(color_hex)

    RX = LX + LW + Inches(0.25)
    RY = LY
    RW = Inches(6.0)
    RH = LH
    steps_box = add_box(
        s,
        RX,
        RY,
        RW,
        RH,
        fill=RGBColor(0xF8, 0xFA, 0xFC),
        line=RGBColor(0xCB, 0xD5, 0xE1),
        line_w=1.0,
        corner=0.02,
    )
    tf = steps_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(260000)
    tf.margin_right = Emu(220000)
    tf.margin_top = Emu(220000)
    tf.margin_bottom = Emu(180000)
    steps = [
        (
            "1.  Embed the query",
            "Marengo via the us.* inference profile — text, image, or text+image — returns a 512-d vector.",
        ),
        (
            "2.  Pull candidate pool",
            "Postgres + pgvector, ORDER BY embedding <=> :q, LIMIT max(80, top_k * 4).",
        ),
        (
            "3.  Frame-snap each clip",
            "For every kind='clip' hit, replace timestamp_sec / thumb_s3_key with the best frame inside [start_sec, end_sec].",
        ),
        (
            "4.  Dedupe within 3s",
            "Walking score-DESC, drop hits within 3s of an earlier survivor in the same s3_key.",
        ),
        (
            "5.  Enrich with presigned URLs",
            "Video URL gets #t=<timestamp_sec>; thumb_s3_key gets its own presigned GET.",
        ),
    ]
    p0 = tf.paragraphs[0]
    for i, (head, body) in enumerate(steps):
        para = p0 if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(6)
        r = para.add_run()
        r.text = head
        r.font.name = "Inter"
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = INK
        para2 = tf.add_paragraph()
        para2.alignment = PP_ALIGN.LEFT
        para2.space_after = Pt(8)
        r2 = para2.add_run()
        r2.text = body
        r2.font.name = "Inter"
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED

    add_text(
        s,
        Inches(0.6),
        Inches(6.7),
        Inches(12.1),
        Inches(0.4),
        "Whole flow lives in app/search.py::search — same pre/post-processing as the local-only scripts/embed/serve.py.",
        size=11,
        color=SUBTLE,
        font="JetBrains Mono",
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 7 — User-facing impact (without vs with frame-snap)
# ---------------------------------------------------------------------------


def slide_user_view(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(s, "What the operator actually sees", kicker="UX impact")

    panels = [
        (
            "WITHOUT frame-snap",
            DANGER,
            RGBColor(0xFE, 0xF2, 0xF2),
            [
                "thumb_s3_key = the clip start frame",
                "presigned URL has #t = 12.00 — start of the 6s clip",
                "operator presses play, then has to scrub forward 3s",
                "result list often shows the same establishing shot",
            ],
            "…/raw-videos/foo.mp4#t=12.00",
            "[ thumbnail = clip start ]",
        ),
        (
            "WITH frame-snap",
            ACCENT,
            RGBColor(0xEC, 0xFD, 0xF5),
            [
                "thumb_s3_key = the matching frame at t = 15.3s",
                "presigned URL has #t = 15.30 — the actual moment",
                "video opens already cued to the matched second",
                "dedupe inside 3s collapses near-duplicate frames",
            ],
            "…/raw-videos/foo.mp4#t=15.30",
            "[ thumbnail = matched frame ]",
        ),
    ]
    px = Inches(0.6)
    py = Inches(1.85)
    pw = Inches(6.0)
    ph = Inches(4.8)
    for title, accent, fill, bullets, url, thumb_label in panels:
        add_box(s, px, py, pw, ph, fill=fill, line=accent, line_w=1.5, corner=0.03)
        add_text(
            s,
            px + Inches(0.3),
            py + Inches(0.25),
            pw - Inches(0.6),
            Inches(0.5),
            title,
            size=14,
            bold=True,
            color=accent,
        )
        by = py + Inches(0.85)
        for b in bullets:
            add_text(
                s,
                px + Inches(0.3),
                by,
                pw - Inches(0.6),
                Inches(0.5),
                "•  " + b,
                size=14,
                color=INK,
            )
            by = by + Inches(0.6)
        thumb_y = py + ph - Inches(1.5)
        add_box(
            s,
            px + Inches(0.3),
            thumb_y,
            Inches(2.0),
            Inches(1.15),
            fill=RGBColor(0xE5, 0xE7, 0xEB),
            line=RGBColor(0x9C, 0xA3, 0xAF),
            line_w=1.0,
            corner=0.04,
        )
        add_text(
            s,
            px + Inches(0.3),
            thumb_y,
            Inches(2.0),
            Inches(1.15),
            thumb_label,
            size=11,
            color=MUTED,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s,
            px + Inches(2.45),
            thumb_y + Inches(0.2),
            pw - Inches(2.75),
            Inches(0.4),
            "presigned URL",
            size=11,
            bold=True,
            color=MUTED,
        )
        add_text(
            s,
            px + Inches(2.45),
            thumb_y + Inches(0.55),
            pw - Inches(2.75),
            Inches(0.5),
            url,
            size=11,
            color=INK,
            font="JetBrains Mono",
        )
        px = px + pw + Inches(0.15)
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 8 — Numbers and knobs
# ---------------------------------------------------------------------------


def slide_numbers(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(s, "Numbers and knobs", kicker="Tuning")

    rows = [
        ("Embed model · Marengo", "twelvelabs.marengo-embed-3-0-v1:0", "us.* inference profile · 512-d L2-norm"),
        ("Text model · Pegasus", "twelvelabs.pegasus-1-2-v1:0", "invoke_model_with_response_stream"),
        ("YOLO models", "pldm-power-line · airpelago-insulator-pole", "weights at s3://.../models/yolo/<name>/v1/best.pt"),
        ("ANN index", "HNSW · m=16 · ef_construction=64", "pgvector on db.t4g.large"),
        ("Clip length", "~6 seconds", "Marengo's native segmentation, async invoke"),
        ("Frame sampling", "1 fps · 720px JPEG · q=4", "ffmpeg inside the frame worker"),
        ("Frame worker (Fargate)", "1024 CPU · 2048 MiB · 30 GiB ephemeral", "8 parallel Bedrock InvokeModel calls"),
        ("Pegasus worker (Fargate)", "1024 CPU · 2048 MiB · 30 GiB ephemeral", "one task per video · prompt preset 'inspector'"),
        ("YOLO worker (Fargate)", "2048 CPU · 4096 MiB · 30 GiB ephemeral", "CPU-only torch · imgsz=640 · conf=0.10 · iou=0.5"),
        ("Lambdas", "start_clip · finalize_clip · start_frame · start_yolo", "256 / 512 / 256 / 256 MiB · pg8000 (pure-Python)"),
        ("Candidate pool", "max(80, top_k × 4)", "covers both clip + frame so frame-snap has room"),
        ("Top-K returned", "1 – 50 (default 10)", "after snap + dedupe"),
        ("Dedupe window", "3.0 seconds", "per s3_key, walking score-DESC"),
        ("Presigned URL TTL", "3600 seconds", "matches existing portal upload flow"),
    ]
    table_x = Inches(0.6)
    table_y = Inches(1.70)
    row_h = Inches(0.32)
    col_widths = [Inches(3.0), Inches(4.4), Inches(4.9)]
    headers = ["Knob", "Value", "Why"]
    cx = table_x
    for h, w in zip(headers, col_widths):
        add_text(s, cx, table_y, w, Inches(0.4), h, size=11, bold=True, color=ACCENT)
        cx = cx + w
    table_y = table_y + Inches(0.4)

    for r_i, row in enumerate(rows):
        if r_i % 2 == 0:
            band_w = Emu(int(col_widths[0]) + int(col_widths[1]) + int(col_widths[2]) + int(Inches(0.1)))
            band = add_box(
                s,
                table_x - Inches(0.05),
                table_y - Inches(0.04),
                band_w,
                row_h,
                fill=RGBColor(0xF8, 0xFA, 0xFC),
                line=RGBColor(0xF8, 0xFA, 0xFC),
                corner=0.03,
            )
            band.shadow.inherit = False
        cx = table_x
        for v, w in zip(row, col_widths):
            mono = (
                v.startswith("twelvelabs.")
                or v.startswith("HNSW")
                or v.startswith("max(")
                or v.startswith("start_clip")
                or v.startswith("1024")
                or v.startswith("2048")
                or v.startswith("8 parallel")
                or v.startswith("pldm-power-line")
            )
            add_text(
                s,
                cx + Inches(0.1),
                table_y + Inches(0.05),
                w - Inches(0.2),
                row_h,
                v,
                size=12,
                color=INK,
                font="JetBrains Mono" if mono else "Inter",
            )
            cx = cx + w
        table_y = table_y + row_h

    add_text(
        s,
        Inches(0.6),
        Inches(6.6),
        Inches(12.1),
        Inches(0.4),
        "Constants live at the top of app/search.py. Defaults are tuned for ~10 videos × hundreds of frames; revisit pool size for >10k clips.",
        size=11,
        color=SUBTLE,
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 9 — Roadmap
# ---------------------------------------------------------------------------


def slide_roadmap(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    add_title_bar(s, "What ships next", kicker="Roadmap")

    phases = [
        (
            "D.1",
            "DONE",
            ACCENT,
            "RDS schema (videos + embeddings) + migrations baked into the portal image. /api/db/health surfaces pgvector status.",
        ),
        (
            "D.2",
            "DONE",
            ACCENT,
            "/api/search/{text,image,text-image} on ECS; Bedrock InvokeModel wired into the task role. Search tab on the portal HUD.",
        ),
        (
            "D.3+4",
            "DONE",
            ACCENT,
            "Marengo embed pipelines: start_clip_embed + finalize_clip_embed Lambdas (VPC, pg8000) for clips; "
            "start_frame_task → frame-embed-worker (Fargate, ffmpeg + 8× parallel Bedrock) for frames.",
        ),
        (
            "D.5",
            "DONE",
            ACCENT,
            "Pegasus per-clip text: finalize_clip_embed now spawns clip-pegasus-worker (Fargate, ffmpeg cuts → "
            "invoke_model_with_response_stream). Descriptions land in clip_descriptions and render inline in the result list.",
        ),
        (
            "D.6",
            "GATED",
            HIGHLIGHT,
            "YOLO instance segmentation: start_yolo_task → yolo-detect-worker (Fargate, ultralytics CPU torch). "
            "Code shipped end-to-end; flips on once the trained best.pt files land at s3://.../models/yolo/.",
        ),
        (
            "V2",
            "NEXT",
            HIGHLIGHT,
            "Embed concatenated detection labels so 'thermal anomaly near a transformer' works without an example image. "
            "Worker DLQs + EB metric alarms; second YOLO model fleet (vegetation / corrosion) when checkpoints land.",
        ),
    ]
    py = Inches(1.75)
    for tag, status, color, body in phases:
        chip = add_box(
            s,
            Inches(0.6),
            py,
            Inches(0.78),
            Inches(0.78),
            fill=color,
            line=color,
            corner=0.18,
        )
        chip.shadow.inherit = False
        add_text(
            s,
            Inches(0.6),
            py,
            Inches(0.78),
            Inches(0.78),
            tag,
            size=15,
            bold=True,
            color=BG,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s,
            Inches(1.55),
            py + Inches(0.02),
            Inches(1.4),
            Inches(0.32),
            status,
            size=10,
            bold=True,
            color=color,
        )
        add_text(
            s,
            Inches(1.55),
            py + Inches(0.30),
            Inches(11.2),
            Inches(0.55),
            body,
            size=13,
            color=INK,
        )
        py = py + Inches(0.85)
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Slide 10 — Close
# ---------------------------------------------------------------------------


def slide_close(prs, idx, total):
    s = add_blank_slide(prs)
    set_slide_background(s, BG)
    band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    add_text(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(0.5),
        "TL;DR",
        size=14,
        bold=True,
        color=ACCENT,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(2.5),
        Inches(11.7),
        Inches(2.4),
        "Clips win recall.\nFrames win the seek bar.\nWe store both and let the search join them.",
        size=42,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.6),
        Inches(11.7),
        Inches(0.6),
        "Demo: search by text, image, or both — and watch the thumbnail land on the actual moment.",
        size=18,
        color=MUTED,
    )
    add_footer(s, idx, total)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_architecture,
        slide_async_pipeline,
        slide_enrichments,
        slide_two_granularities,
        slide_frame_snap,
        slide_search_algo,
        slide_user_view,
        slide_numbers,
        slide_roadmap,
        slide_close,
    ]
    total = len(builders)
    for i, b in enumerate(builders, start=1):
        b(prs, i, total)

    prs.save(OUT)
    print(f"wrote {OUT.relative_to(DOCS.parent)}")


if __name__ == "__main__":
    main()
